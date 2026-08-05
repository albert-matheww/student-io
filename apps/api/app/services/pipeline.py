"""The AI processing pipeline: OCR -> speech-to-text -> text extraction ->
chunking -> semantic embeddings — the flow described in the product spec.
Runs as a FastAPI background task immediately after a resource is uploaded,
moving `Resource.status` through uploaded -> processing -> processed/failed.
"""

import io

from app.core.config import get_settings
from app.core.database import SessionLocal
from app.models import Resource, ResourceChunk, ResourceStatus, ResourceType
from app.services import ai
from app.services.storage import read_resource_file

settings = get_settings()

CHUNK_SIZE = 1200
CHUNK_OVERLAP = 150


def process_resource(resource_id: str) -> None:
    """Background-task entry point. Owns its own DB session since it runs
    after the request/response cycle that created the original session."""
    db = SessionLocal()
    try:
        resource = db.query(Resource).filter(Resource.id == resource_id).first()
        if resource is None:
            return

        resource.status = ResourceStatus.processing
        resource.processing_stage = "extracting_text"
        db.commit()

        try:
            content = read_resource_file(resource.storage_path)
            text = _extract_text(resource.resource_type, resource.filename, content)
        except Exception as exc:  # noqa: BLE001 — any extraction failure is reported to the user, not just logged
            resource.status = ResourceStatus.failed
            resource.error_message = str(exc)[:500]
            db.commit()
            return

        if not text.strip():
            resource.status = ResourceStatus.processed
            resource.processing_stage = "no_extractable_text"
            db.commit()
            return

        resource.processing_stage = "chunking"
        db.commit()
        chunks = _chunk_text(text)

        resource.processing_stage = "embedding"
        db.commit()
        try:
            chunks_with_embeddings = []
            for index, chunk in enumerate(chunks):
                embedding = ai.embed_text(chunk) if ai.active_provider() else None
                chunks_with_embeddings.append(
                    ResourceChunk(resource_id=resource.id, content=chunk, chunk_index=index, embedding=embedding)
                )
        except Exception as exc:  # noqa: BLE001
            resource.status = ResourceStatus.failed
            resource.error_message = str(exc)[:500]
            db.commit()
            return
        for chunk_row in chunks_with_embeddings:
            db.add(chunk_row)

        resource.status = ResourceStatus.processed
        resource.processing_stage = "done"
        db.commit()
    finally:
        db.close()


def _extract_text(resource_type: ResourceType, filename: str, content: bytes) -> str:
    if resource_type == ResourceType.pdf:
        return _extract_pdf(content)
    if resource_type == ResourceType.docx:
        return _extract_docx(content)
    if resource_type == ResourceType.ppt:
        return _extract_pptx(content)
    if resource_type == ResourceType.txt:
        return content.decode("utf-8", errors="ignore")
    if resource_type == ResourceType.image:
        return _transcribe_image(content, filename)
    if resource_type in (ResourceType.audio, ResourceType.video):
        return _transcribe_audio(content, filename)
    return ""


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def _transcribe_image(content: bytes, filename: str) -> str:
    """OCR + diagram description via a vision-capable model."""
    if not ai.active_provider():
        return ""
    return ai.transcribe_image(content, filename)


def _transcribe_audio(content: bytes, filename: str) -> str:
    if not ai.active_provider():
        return ""
    return ai.transcribe_audio(content, filename)


def _chunk_text(text: str) -> list[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        chunks.append(text[start:end])
        start = max(end - CHUNK_OVERLAP, start + 1)
    return [c.strip() for c in chunks if c.strip()]
